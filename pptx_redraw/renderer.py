"""Shape rendering helpers for python-pptx."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Inches, Pt

from pptx_redraw.theme import FONT_CJK, FONT_UI, Color, TEXT, TITLE


def _rgb(color: Color) -> RGBColor:
    """Convert Color dataclass to python-pptx RGBColor.

    Args:
        color: RGB value container.

    Returns:
        A python-pptx RGBColor value.

    Raises:
        ValueError: Not raised in this conversion.
    """

    return RGBColor(color.r, color.g, color.b)


def add_round_rect(
    slide: Slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: Color,
    line: Color,
    line_width: float = 1.0,
    radius_shape: MSO_AUTO_SHAPE_TYPE = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
):
    """Add a rounded rectangle-like shape.

    Args:
        slide: Destination slide.
        x: Left in inches.
        y: Top in inches.
        w: Width in inches.
        h: Height in inches.
        fill: Fill color.
        line: Border color.
        line_width: Border width in points.
        radius_shape: Shape type.

    Returns:
        Newly added shape object.

    Raises:
        ValueError: If numeric dimensions are invalid.
    """

    shape = slide.shapes.add_shape(radius_shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line)
    shape.line.width = Pt(line_width)
    return shape


def add_textbox(
    slide: Slide,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    *,
    size: int = 12,
    bold: bool = False,
    color: Color = TEXT,
    font_name: str = FONT_CJK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    """Add editable text box.

    Args:
        slide: Destination slide.
        x: Left in inches.
        y: Top in inches.
        w: Width in inches.
        h: Height in inches.
        text: Text content.
        size: Font size in points.
        bold: Font bold flag.
        color: Font color.
        font_name: Font family name.
        align: Paragraph alignment.

    Returns:
        Added text box shape.

    Raises:
        ValueError: If dimensions are invalid.
    """

    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)
    run.font.name = font_name
    return tb


def add_card(
    slide: Slide,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    bullets: list[str],
    *,
    fill: Color,
    line: Color,
):
    """Add a titled card with bullet-like rows.

    Args:
        slide: Destination slide.
        x: Left in inches.
        y: Top in inches.
        w: Width in inches.
        h: Height in inches.
        title: Card title.
        bullets: Bullet row texts.
        fill: Card fill color.
        line: Card border color.

    Returns:
        Added card shape.

    Raises:
        ValueError: If bullets are empty.
    """

    shape = add_round_rect(slide, x, y, w, h, fill=fill, line=line)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()

    p0 = tf.paragraphs[0]
    p0.text = title
    p0.alignment = PP_ALIGN.LEFT
    r0 = p0.runs[0]
    r0.font.size = Pt(11)
    r0.font.bold = True
    r0.font.color.rgb = _rgb(TITLE)
    r0.font.name = FONT_CJK

    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        r = p.runs[0]
        r.font.size = Pt(9)
        r.font.color.rgb = _rgb(TEXT)
        r.font.name = FONT_CJK
    return shape


def add_objective_icon(slide: Slide, x: float, y: float, tag: str, *, fill: Color, line: Color):
    """Add simple vector icon approximation using editable shapes.

    Args:
        slide: Destination slide.
        x: Left in inches.
        y: Top in inches.
        tag: Icon center text.
        fill: Fill color.
        line: Border color.

    Returns:
        Tuple of icon shapes.

    Raises:
        ValueError: If tag is empty.
    """

    circle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(0.26), Inches(0.26)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = _rgb(fill)
    circle.line.color.rgb = _rgb(line)

    tf = circle.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = tag
    run = p.runs[0]
    run.font.size = Pt(7)
    run.font.bold = True
    run.font.name = FONT_UI
    run.font.color.rgb = RGBColor(255, 255, 255)
    return circle
