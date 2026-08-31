"""Export pipeline for editable AegisOps architecture PPTX."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.util import Inches, Pt

from pptx_redraw.models import ObjectiveItem, RoadmapQuarter, SectionCard, SlideContent
from pptx_redraw.renderer import add_card, add_objective_icon, add_round_rect, add_textbox
from pptx_redraw.theme import (
    ACCENT,
    ACCENT_SOFT,
    BG,
    BORDER,
    GOOD,
    MUTED,
    ORANGE,
    PANEL,
    PURPLE,
    TEAL,
    TEXT,
    TITLE,
)


def _rgb(color) -> RGBColor:
    """Convert theme color to RGBColor.

    Args:
        color: Theme color object with r/g/b channels.

    Returns:
        Converted RGBColor.

    Raises:
        ValueError: Not raised by this converter.
    """

    return RGBColor(color.r, color.g, color.b)


def build_content() -> SlideContent:
    """Build strongly typed slide content model.

    Args:
        None.

    Returns:
        SlideContent instance used by renderer.

    Raises:
        ValueError: If model validation fails.
    """

    return SlideContent(
        title="AegisOps 最小可實作架構 (MVP) & 分期實施計畫",
        subtitle="企業級 AgentOps 架構藍圖（完全可編輯重繪版）",
        badge="目標：先可觀測，再可評估，最後可治理與自動化",
        control_sections=[
            SectionCard(
                title="1. 核心控制層 (Control Plane)",
                bullets=[
                    "Agent Registry",
                    "Prompt Registry",
                    "Model Registry",
                    "Tool / MCP Registry",
                    "Policy & Guardrails",
                    "RBAC / IAM",
                ],
            ),
            SectionCard(
                title="2. 觀測與成本 (Observability & Cost)",
                bullets=[
                    "Trace Collector (OpenTelemetry)",
                    "Token / Trace Store (Traces / Spans)",
                    "Metrics Store (Metrics)",
                    "Cost Tracker (Token & $)",
                ],
            ),
            SectionCard(
                title="3. 評估與品質 (Evaluation & Quality)",
                bullets=[
                    "Evaluation Engine (LLM Judge / Rule)",
                    "Quality Dashboard (品質指標)",
                    "Dataset Store (Eval Dataset)",
                ],
            ),
            SectionCard(
                title="4. 事件・回饋與風險 (Incident / Feedback / Risk)",
                bullets=[
                    "Incident Center (事件管理)",
                    "Feedback Collector (使用者回饋)",
                    "Risk Detector (風險檢測)",
                ],
            ),
            SectionCard(
                title="5. 資料儲存層",
                bullets=[
                    "PostgreSQL (Meta / 統計資料)",
                    "S3 / Blob Storage (檔案 / Prompt / Dataset / Log)",
                    "Vector DB (向量 / 知識庫)",
                    "Cache (Redis)",
                ],
            ),
            SectionCard(
                title="6. 整合與交付層",
                bullets=[
                    "API Gateway (內部 API)",
                    "Webhook / Events (事件通知)",
                    "CI/CD Hook (部署觸發)",
                    "Notification (Email / Slack)",
                ],
            ),
        ],
        objectives=[
            ObjectiveItem(name="Token / Trace", description="追蹤每次呼叫與上下文鏈路"),
            ObjectiveItem(name="Evaluation", description="LLM Judge 與規則化評測"),
            ObjectiveItem(name="Cost", description="Token 與 $ 成本可見化"),
            ObjectiveItem(name="Quality", description="品質指標與趨勢監控"),
            ObjectiveItem(name="Incident", description="事件管理與回應流程"),
            ObjectiveItem(name="Feedback", description="回饋蒐集、分類與閉環"),
            ObjectiveItem(name="Risk", description="風險檢測與告警機制"),
        ],
        roadmap=[
            RoadmapQuarter(
                quarter="Q1：可觀測 (See)",
                summary="建立可觀測與成本基礎能力",
                bullets=[
                    "接入 OpenTelemetry、Trace Collector 與 Trace Store",
                    "Token / Cost 追蹤",
                    "基礎 Dashboard (Trace / Cost)",
                ],
            ),
            RoadmapQuarter(
                quarter="Q2：評估與品質 (Measure)",
                summary="建立評估與品質衡量指標",
                bullets=[
                    "Evaluation Engine (LLM Judge / Rule)",
                    "Dataset 管理",
                    "Quality Dashboard (品質指標)、Evaluation API",
                ],
            ),
            RoadmapQuarter(
                quarter="Q3：事件與回饋 (React)",
                summary="建立事件與使用者回饋閉環機制",
                bullets=[
                    "Incident Center (事件管理)",
                    "Feedback Collector (回饋收集)",
                    "Risk Detector (風險檢測)",
                    "告警與通知機制",
                ],
            ),
            RoadmapQuarter(
                quarter="Q4：治理與自動化 (Control)",
                summary="強化治理、發布與自動化能力",
                bullets=[
                    "Agent / Prompt / Tool 全量管理",
                    "Policy & Guardrails 強化",
                    "CI/CD Integration",
                    "報表與稽核 (Audit)",
                ],
            ),
        ],
    )


def generate_pptx(output_path: Path) -> Path:
    """Generate single-slide editable PPTX.

    Args:
        output_path: Destination file path.

    Returns:
        Resolved output path.

    Raises:
        OSError: If save fails.
    """

    content = build_content()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(BG)

    add_textbox(slide, 0.3, 0.1, 7.8, 0.35, content.title, size=20, bold=True, color=TITLE)
    add_textbox(slide, 0.3, 0.42, 5.8, 0.22, content.subtitle, size=11, color=MUTED)
    badge = add_round_rect(slide, 9.1, 0.12, 4.0, 0.42, fill=ACCENT_SOFT, line=ACCENT)
    badge.text_frame.text = content.badge
    badge.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
    badge.text_frame.paragraphs[0].runs[0].font.bold = True

    add_round_rect(slide, 0.3, 0.75, 3.15, 1.65, fill=PANEL, line=BORDER)
    add_textbox(slide, 0.42, 0.82, 2.9, 0.2, "Agent Runtime (執行層)", size=11, bold=True)
    add_textbox(
        slide,
        0.45,
        1.05,
        2.85,
        1.2,
        "• LangGraph\n• OpenAI Agents SDK\n• Dify / 其他 Agent 框架\n• MCP / Tool / API\n\n外部系統 / 資料源\n• 企業資料庫\n• 向量資料庫\n• 外部 API / 服務\n• 檔案 / 文件 / Web",
        size=9,
    )

    add_round_rect(slide, 3.7, 0.75, 5.85, 0.58, fill=PANEL, line=BORDER)
    add_textbox(slide, 3.82, 0.82, 1.3, 0.2, "使用者 / 存取入口", size=10, bold=True)
    entry_items = ["Web Portal", "內部人員", "開發人員", "營運 / 維運", "管理者"]
    for idx, item in enumerate(entry_items):
        x = 5.25 + idx * 0.84
        box = add_round_rect(slide, x, 0.91, 0.77, 0.28, fill=ACCENT_SOFT, line=BORDER)
        box.text_frame.text = item
        run = box.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(7)

    main = add_round_rect(slide, 3.7, 1.4, 5.85, 3.25, fill=PANEL, line=ACCENT, line_width=1.5)
    main.text_frame.text = "AegisOps 最小可實作架構 (MVP)"
    mr = main.text_frame.paragraphs[0].runs[0]
    mr.font.bold = True
    mr.font.size = Pt(9)

    card_w, card_h = 2.75, 0.93
    start_x, start_y = 3.88, 1.82
    for i, sec in enumerate(content.control_sections):
        row = i // 2
        col = i % 2
        add_card(
            slide,
            start_x + col * 2.95,
            start_y + row * 1.01,
            card_w,
            card_h,
            sec.title,
            sec.bullets,
            fill=ACCENT_SOFT if i % 2 == 0 else PANEL,
            line=BORDER,
        )

    add_round_rect(slide, 9.75, 0.75, 3.25, 3.9, fill=PANEL, line=BORDER)
    add_textbox(slide, 9.88, 0.82, 2.9, 0.2, "7 大著要目的 (MVP 聚焦)", size=11, bold=True)
    icon_tags = ["T", "E", "$", "Q", "I", "F", "R"]
    icon_fills = [ACCENT, PURPLE, GOOD, TEAL, ORANGE, ACCENT, PURPLE]
    for idx, item in enumerate(content.objectives):
        y = 1.1 + idx * 0.47
        add_objective_icon(slide, 9.92, y, icon_tags[idx], fill=icon_fills[idx], line=icon_fills[idx])
        add_textbox(slide, 10.23, y - 0.01, 1.0, 0.15, item.name, size=9, bold=True)
        add_textbox(slide, 11.2, y - 0.01, 1.7, 0.22, item.description, size=8, color=MUTED)

    add_round_rect(slide, 0.3, 4.78, 12.7, 1.62, fill=PANEL, line=BORDER)
    add_textbox(slide, 0.42, 4.84, 3.2, 0.22, "分期實施計畫 (最小可行路線圖)", size=11, bold=True)

    q_colors = [ACCENT_SOFT, PANEL, ACCENT_SOFT, PANEL]
    for i, quarter in enumerate(content.roadmap):
        x = 0.42 + i * 3.12
        card = add_round_rect(slide, x, 5.08, 2.95, 1.2, fill=q_colors[i], line=BORDER)
        tf = card.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        p0.text = quarter.quarter
        p0.runs[0].font.bold = True
        p0.runs[0].font.size = Pt(7)
        p1 = tf.add_paragraph()
        p1.text = quarter.summary
        p1.runs[0].font.size = Pt(6)
        for bullet in quarter.bullets:
            p = tf.add_paragraph()
            p.text = f"• {bullet}"
            p.runs[0].font.size = Pt(5.5)

        if i < 3:
            arr = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                Inches(x + 2.98),
                Inches(5.54),
                Inches(0.12),
                Inches(0.16),
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = _rgb(ACCENT)
            arr.line.fill.background()

    add_round_rect(slide, 0.3, 6.5, 12.7, 0.78, fill=PANEL, line=BORDER)
    add_textbox(slide, 0.42, 6.56, 2.8, 0.2, "技術基礎建議 (可彈性替換)", size=10, bold=True)
    add_textbox(
        slide,
        2.95,
        6.56,
        8.4,
        0.22,
        "OpenTelemetry｜PostgreSQL｜Redis｜S3 / Blob｜Chroma / pgvector｜FastAPI｜Langfuse｜Phoenix｜Grafana｜Prometheus｜Docker / K8s",
        size=8,
    )
    add_textbox(slide, 11.5, 6.56, 1.5, 0.2, "備註：可依企業規範替換", size=8, color=MUTED)

    c1 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(3.42), Inches(1.55), Inches(3.68), Inches(1.55)
    )
    c1.line.color.rgb = _rgb(ACCENT)
    c1.line.width = Pt(0.75)

    c2 = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(9.56), Inches(2.6), Inches(9.73), Inches(2.6)
    )
    c2.line.color.rgb = _rgb(ACCENT)
    c2.line.width = Pt(0.75)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path.resolve()
